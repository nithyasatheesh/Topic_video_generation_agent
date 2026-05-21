import streamlit as st
from openai import OpenAI

from PIL import (
    Image,
    ImageDraw,
    ImageFont
)

from moviepy.editor import (
    ImageClip,
    AudioFileClip,
    concatenate_videoclips
)

from pptx import Presentation

import subprocess
import os
import json
import base64
import shutil

from io import BytesIO

st.set_page_config(
    page_title="AI PPT / Topic Video Generator",
    layout="wide"
)

st.title("🎬 PPT / Topic → Video Generator")

client = OpenAI(
    api_key=st.secrets["OPENAI_API_KEY"]
)

TEMP_DIR = "temp"

WIDTH = 1280
HEIGHT = 720

os.makedirs(
    TEMP_DIR,
    exist_ok=True
)


def cleanup():

    if os.path.exists(TEMP_DIR):

        shutil.rmtree(TEMP_DIR)

    os.makedirs(
        TEMP_DIR,
        exist_ok=True
    )


def font(size):

    try:

        return ImageFont.truetype(
            "DejaVuSans.ttf",
            size
        )

    except:

        return ImageFont.load_default()


####################################################
# TOPIC MODE
####################################################

def generate_content(topic):

    prompt = f"""
Create educational content about:

{topic}

Return JSON:

{{
"title":"",
"slides":[
{{
"heading":"",
"content":"",
"image_prompt":""
}}
]
}}

Generate exactly 5 slides.
Maximum 25 words content.
"""

    response = client.chat.completions.create(
        model="gpt-4.1",
        messages=[
            {
                "role":"user",
                "content":prompt
            }
        ],
        response_format={
            "type":"json_object"
        }
    )

    return json.loads(
        response.choices[0].message.content
    )


def generate_image(prompt,index):

    image = client.images.generate(
        model="gpt-image-1",
        prompt=prompt,
        size="1536x1024",
        quality="high"
    )

    data = base64.b64decode(
        image.data[0].b64_json
    )

    img = Image.open(
        BytesIO(data)
    ).convert("RGB")

    path = f"{TEMP_DIR}/image_{index}.png"

    img.save(path)

    return path


def wrap_text(draw,text,font,width):

    words=text.split()

    lines=[]

    current=""

    for word in words:

        test=current+" "+word

        size=draw.textbbox(
            (0,0),
            test,
            font=font
        )[2]

        if size>width:

            lines.append(current)

            current=word

        else:

            current=test

    lines.append(current)

    return "\n".join(lines)


def create_slide(
    heading,
    content,
    image_path,
    idx
):

    slide=Image.new(
        "RGB",
        (WIDTH,HEIGHT),
        "white"
    )

    img=Image.open(
        image_path
    )

    img.thumbnail(
        (500,500)
    )

    slide.paste(
        img,
        (700,100)
    )

    draw=ImageDraw.Draw(
        slide
    )

    title=font(52)

    body=font(30)

    draw.text(
        (60,80),
        heading,
        fill="black",
        font=title
    )

    draw.multiline_text(
        (60,180),
        wrap_text(
            draw,
            content,
            body,
            550
        ),
        fill="black",
        font=body,
        spacing=12
    )

    path=f"{TEMP_DIR}/slide_{idx}.png"

    slide.save(path)

    return path


####################################################
# PPT MODE
####################################################

def ppt_to_images(ppt_file):

    ppt_path=f"{TEMP_DIR}/input.pptx"

    with open(
        ppt_path,
        "wb"
    ) as f:

        f.write(
            ppt_file.read()
        )

    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        ppt_path,
        "--outdir",
        TEMP_DIR
    ])

    pdf=f"{TEMP_DIR}/input.pdf"

    subprocess.run([
        "pdftoppm",
        "-png",
        pdf,
        f"{TEMP_DIR}/slide"
    ])

    files=sorted([
        os.path.join(
            TEMP_DIR,
            x
        )
        for x in os.listdir(
            TEMP_DIR
        )
        if x.startswith("slide-")
    ])

    return files


def extract_slide_text(ppt):

    prs=Presentation(
        ppt
    )

    output=[]

    for slide in prs.slides:

        text=[]

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                text.append(
                    shape.text
                )

        output.append(
            " ".join(text)
        )

    return output


####################################################
# AUDIO
####################################################

def narration(
    text,
    idx
):

    speech=client.audio.speech.create(
        model="gpt-4o-mini-tts",
        voice="alloy",
        input=text
    )

    path=f"{TEMP_DIR}/audio_{idx}.mp3"

    speech.stream_to_file(
        path
    )

    return path


####################################################
# VIDEO
####################################################

def build_video(
    images,
    audio
):

    clips=[]

    for img,aud in zip(
        images,
        audio
    ):

        sound=AudioFileClip(
            aud
        )

        clip=(
            ImageClip(img)
            .set_duration(
                sound.duration
            )
            .set_audio(sound)
        )

        clips.append(
            clip
        )

    final=concatenate_videoclips(
        clips,
        method="compose"
    )

    output=f"{TEMP_DIR}/video.mp4"

    final.write_videofile(
        output,
        codec="libx264",
        audio_codec="aac",
        fps=24,
        logger=None
    )

    return output


mode=st.radio(
    "Mode",
    [
        "Topic",
        "PPT Upload"
    ]
)

cleanup()

slides=[]
audio=[]

if mode=="Topic":

    topic=st.text_input(
        "Topic"
    )

    if st.button(
        "Generate"
    ):

        data=generate_content(
            topic
        )

        for idx,s in enumerate(
            data["slides"]
        ):

            img=generate_image(
                s["image_prompt"],
                idx
            )

            slide=create_slide(
                s["heading"],
                s["content"],
                img,
                idx
            )

            aud=narration(
                s["content"],
                idx
            )

            slides.append(
                slide
            )

            audio.append(
                aud
            )

elif mode=="PPT Upload":

    ppt=st.file_uploader(
        "Upload PPT",
        type=["pptx"]
    )

    if ppt:

        ppt_path=f"{TEMP_DIR}/ppt.pptx"

        with open(
            ppt_path,
            "wb"
        ) as f:

            f.write(
                ppt.read()
            )

        slides=ppt_to_images(
            open(
                ppt_path,
                "rb"
            )
        )

        texts=extract_slide_text(
            ppt_path
        )

        for idx,text in enumerate(
            texts
        ):

            audio.append(
                narration(
                    text,
                    idx
                )
            )

if slides and audio:

    with st.spinner(
        "Creating video..."
    ):

        video=build_video(
            slides,
            audio
        )

    st.video(video)

    with open(
        video,
        "rb"
    ) as f:

        st.download_button(
            "Download Video",
            f,
            file_name="video.mp4"
        )
